import streamlit as st
import json
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO

# ==========================================
# 1. AI 結構化核心 (Groq Llama-3)
# ==========================================
def structure_content_with_ai(topic, outline, content, api_key):
    client = Groq(api_key=api_key)
    
    # 定義輸出的 JSON 格式
    json_structure = """
    {
        "title_slide": { "title": "Main Title", "subtitle": "Subtitle" },
        "content_slides": [
            { "title": "Slide Title", "points": ["Point 1", "Point 2", "Point 3"] }
        ]
    }
    """
    
    # Prompt 設計：要求 AI 根據使用者的內容進行整理，而非無中生有
    prompt = f"""
    You are a professional presentation editor.
    Your task is to organize the user's raw input into a structured presentation format.
    
    User Input:
    - Topic: {topic}
    - Outline: {outline}
    - Raw Content: {content}
    
    Instructions:
    1. Create a Title Slide based on the Topic.
    2. Create Content Slides based on the Outline and Raw Content.
    3. Summarize the Raw Content into concise bullet points (3-5 points per slide).
    4. Language: Traditional Chinese (繁體中文).
    5. Output strictly valid JSON following this structure: {json_structure}
    
    Do not add extra conversational text. Output JSON only.
    """
    
    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "You are a JSON-only output assistant."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5, # 降低隨機性，更忠於原文
            response_format={"type": "json_object"}
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        st.error(f"AI 整理失敗: {e}")
        return None

# ==========================================
# 2. PPT 生成核心 (Python-pptx)
# ==========================================
def create_styled_ppt(data, style_config):
    prs = Presentation()
    
    # 解構風格設定
    bg_rgb = style_config["bg_color"]
    title_rgb = style_config["title_color"]
    text_rgb = style_config["text_color"]
    accent_rgb = style_config["accent_color"]
    font_name = style_config["font"]
    
    # --- 1. 封面頁 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白版型
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_rgb
    
    # 封面裝飾邏輯
    if style_config["style_type"] == "geometric":
        # 幾何風格：加入色塊
        shape = slide.shapes.add_shape(1, 0, Inches(6), Inches(10), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
    elif style_config["style_type"] == "organic":
        # 圓潤風格：加入圓形
        oval = slide.shapes.add_shape(9, Inches(6), Inches(-2), Inches(6), Inches(6))
        oval.fill.solid()
        oval.fill.fore_color.rgb = accent_rgb
        oval.line.fill.background()

    # 封面文字
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p = title_box.text_frame.add_paragraph()
    p.text = data["title_slide"]["title"]
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = title_rgb
    p.font.name = font_name
    
    sub = title_box.text_frame.add_paragraph()
    sub.text = data["title_slide"]["subtitle"]
    sub.font.size = Pt(24)
    sub.font.color.rgb = text_rgb
    sub.font.name = font_name

    # --- 2. 內文頁 ---
    for slide_data in data["content_slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_rgb
        
        # 內文頁標題
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tp = t_box.text_frame.add_paragraph()
        tp.text = slide_data["title"]
        tp.font.bold = True
        tp.font.size = Pt(32)
        tp.font.color.rgb = title_rgb
        tp.font.name = font_name
        
        # 標題裝飾線
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(9), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_rgb
        line.line.fill.background()

        # 內容列表
        c_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
        tf = c_box.text_frame
        tf.word_wrap = True
        
        for point in slide_data["points"]:
            cp = tf.add_paragraph()
            cp.text = "• " + point
            cp.font.size = Pt(20)
            cp.font.color.rgb = text_rgb
            cp.font.name = font_name
            cp.space_after = Pt(14)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ==========================================
# 3. Streamlit UI
# ==========================================
st.set_page_config(page_title="Custom Content PPT Generator", layout="wide")

st.title("📝 自訂內容 PPT 生成器")
st.markdown("貼上您的文字內容，AI 將為您整理結構並套用兩種設計風格。")

# Sidebar: 設定 API
with st.sidebar:
    st.header("🔑 設定")
    api_key = st.text_input("Groq API Key", type="password")
    st.caption("需要 Groq API Key 來進行內容結構化")

# Input Form
with st.form("content_form"):
    col1, col2 = st.columns([1, 2])
    
    with col1:
        topic_input = st.text_input("1. 簡報主題", "2024 年度行銷計畫")
        outline_input = st.text_area("2. 大綱 (或是章節)", "市場分析\n策略規劃\n預算分配", height=150)
    
    with col2:
        content_input = st.text_area("3. 詳細內容 (直接貼上文章或草稿)", 
                                   "在市場分析部分，我們發現競爭對手A的市佔率下降了5%。\n"
                                   "我們的策略是專注於社群媒體行銷，特別是 Instagram Reels。\n"
                                   "預算方面，我們預計投入 30% 在廣告投放，20% 在 KOL 合作...", 
                                   height=230)
    
    submit_btn = st.form_submit_button("🚀 開始製作 PPT")

# Processing
if submit_btn:
    if not api_key:
        st.error("請先在左側輸入 Groq API Key！")
    else:
        with st.spinner("🤖 AI 正在閱讀您的內容並進行排版..."):
            # 1. 呼叫 AI 整理內容
            ppt_structure = structure_content_with_ai(topic_input, outline_input, content_input, api_key)
            
            if ppt_structure:
                st.success("內容結構化完成！正在生成檔案...")
                
                # 預覽整理好的結構
                with st.expander("查看 AI 整理後的結構"):
                    st.json(ppt_structure)
                
                # 定義兩種風格
                style_a = {
                    "style_type": "geometric",
                    "bg_color": RGBColor(255, 255, 255),    # 白底
                    "title_color": RGBColor(44, 62, 80),    # 深藍灰
                    "text_color": RGBColor(52, 73, 94),     # 灰藍
                    "accent_color": RGBColor(231, 76, 60),  # 紅色點綴
                    "font": "Microsoft JhengHei"
                }
                
                style_b = {
                    "style_type": "organic",
                    "bg_color": RGBColor(33, 33, 33),       # 深灰底
                    "title_color": RGBColor(241, 196, 15),  # 金黃
                    "text_color": RGBColor(236, 240, 241),  # 淺灰
                    "accent_color": RGBColor(46, 204, 113), # 翠綠點綴
                    "font": "Microsoft JhengHei"
                }

                # 生成兩個檔案
                file_a = create_styled_ppt(ppt_structure, style_a)
                file_b = create_styled_ppt(ppt_structure, style_b)
                
                # 下載區域
                st.markdown("### 🎉 您的 PPT 已準備好")
                d_col1, d_col2 = st.columns(2)
                
                with d_col1:
                    st.image("https://placehold.co/400x200/FFFFFF/2C3E50/png?text=Modern+Business", caption="Style A: 現代商務 (明亮)")
                    st.download_button("下載 Style A", file_a, "presentation_style_a.pptx")
                    
                with d_col2:
                    st.image("https://placehold.co/400x200/212121/F1C40F/png?text=Creative+Dark", caption="Style B: 創意暗色 (高對比)")
                    st.download_button("下載 Style B", file_b, "presentation_style_b.pptx")